"""
Generate an EDITABLE PowerPoint replica of the
"AI's impact across the Tech Org: Capability View" slide.

Every capability is its own native rectangle + text box, so each one can be
moved, recolored, or renamed individually in PowerPoint. Nothing is a flattened
image.

Edit the DATA section below to change text, the COLORS map to change shades, and
the IMPACT overrides dict to change a single box's color.

Run:  python generate_capability_view.py
Out:  AI_Capability_View.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# COLORS  (three-tier AI-impact scale from the legend, easily editable)
# ---------------------------------------------------------------------------
POWERED   = RGBColor(0x4C, 0xA7, 0x2C)   # AI Powered   80/20  (dark green)
AUGMENTED = RGBColor(0xA9, 0xD1, 0x8E)   # AI Augmented 50/50  (medium green)
ASSISTED  = RGBColor(0xE2, 0xEF, 0xDA)   # AI Assisted  20/80  (light green)
HEADER_BG = RGBColor(0x00, 0x00, 0x00)   # black header bars
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BOX_LINE  = RGBColor(0xFF, 0xFF, 0xFF)   # thin white gaps between boxes
TEXT_DK   = RGBColor(0x1A, 0x1A, 0x1A)

# A varied-but-deterministic mix so the board reads like the original
# (mostly augmented/assisted with scattered powered). Colors are APPROXIMATE
# and meant to be edited per box in PowerPoint.
PATTERN = [AUGMENTED, ASSISTED, AUGMENTED, POWERED, ASSISTED,
           AUGMENTED, AUGMENTED, ASSISTED, POWERED, AUGMENTED,
           ASSISTED, AUGMENTED, POWERED, ASSISTED, AUGMENTED]

def color_for(i):
    return PATTERN[i % len(PATTERN)]

# ---------------------------------------------------------------------------
# DATA  (panel title -> number of box columns -> list of capabilities)
# ---------------------------------------------------------------------------
PRODUCT_PANEL = {
    "title": "Product & Solution Development",
    "ncols": 2,
    "items": [
        "Research (Internal & External)", "Product Portfolio Mgmt",
        "Business/Functional Requirements Mgmt", "Product Pricing/Financial Mgmt",
        "Product KPI Mgmt", "Product Roadmapping",
        "Project/Program Mgmt", "Agile Planning & Mgmt",
        "Agile Services & Standards", "Product & Service Design",
        "Product Development & Engineering", "Product Build",
        "Test Planning & Mgmt", "System & Integration Testing",
        "Load & Performance Testing", "User Acceptance Testing",
        "Regression & End-to-End Testing", "Product Configuration",
        "Service Catalog Mgmt", "Capacity & Performance Mgmt",
        "Availability Mgmt", "Service Continuity/Disaster Recovery",
        "Change Control", "Knowledge Mgmt",
        "Release Mgmt", "Monitoring & Event Mgmt",
        "Incident Mgmt", "Request Fulfillment",
        "Access Mgmt", "Problem Mgmt",
    ],
}

TBM = {
    "title": "Tech Business Management", "ncols": 3,
    "items": [
        "Portfolio Mgmt/Prioritization/Reporting", "Portfolio Governance", "Demand Intake & Resource Mgmt",
        "Tech Investment Mgmt", "Business Case Development", "Resource Mgmt",
        "Benefits Tracking", "Financial Planning & Mgmt", "IT Budgeting & Forecasting",
        "IT Capital & Expense Mgmt", "IT Financial Reporting", "Fixed Asset Mgmt",
        "Service Cost & Visibility", "IT Project/Product Costing",
    ],
}

TSI = {
    "title": "Tech Strategy & Innovation", "ncols": 2,
    "items": [
        "Technology Strategy Development", "Innovation Intent & Strategy",
        "Sensing & Scanning", "Innovation Portfolio Mgmt",
        "Prototyping & Proof of Concept Development", "Incubation Platforms",
        "Automation Platform Mgmt", "AI Platform Mgmt",
        "Automation & Regression", "Performance & Volumetrics",
    ],
}

ISM = {
    "title": "Information Security Management", "ncols": 4,
    "items": [
        "Asset Mgmt", "Strategy Dev & Mgmt", "Business Engagement", "Security Policies & Standards",
        "Security Metrics & Reporting", "Regulatory Reqs & Compliance", "Security Communications", "Security Risk Assessment",
        "Secure Mergers & Acquisitions", "Risk Mgmt", "Issues Mgmt", "Operations Tech Security",
        "Third Party Risk Mgmt", "Identity & Access Mgmt", "Data Security Governance", "Data ID & Classification",
        "Data Encryption", "Data Loss Prevention", "Certification & Key Mgmt", "Security Awareness & Training",
        "Application Security", "Cloud Security", "Vulnerability Mgmt", "Configuration Mgmt",
        "Endpoint Security", "Network Security", "Security Incident & Event Monitoring", "Security Operations Center",
        "Threat Intel/Threat Hunting", "Penetration Testing", "Adversarial Simulation", "Crisis Response",
        "Incident Response Plan", "Disaster Recovery", "Business Continuity",
    ],
}

EA = {
    "title": "Enterprise Architecture", "ncols": 3,
    "items": [
        "Enterprise Architecture Repositories", "Solution/Product Architecture Mgmt", "Business Architecture",
        "Security Architecture", "Infrastructure Architecture", "App & Integration Architecture",
        "Platform Architecture", "Portfolio Architecture Mgmt", "Architecture Governance",
        "Network Architecture", "Data Architecture", "Infrastructure Architecture",
    ],
}

IE = {
    "title": "Infrastructure Engineering", "ncols": 2,
    "items": [
        "Server Mgmt (Design, Build, Deploy, Support)", "Storage Mgmt (Design, Build, Deploy, Support)",
        "Database Administration", "Data Center Operations",
        "Infrastructure Reliability Engineering", "Disaster Recovery",
    ],
}

DA = {
    "title": "Data & Analytics", "ncols": 3,
    "items": [
        "Standardized Reporting & Visualization", "Reports & Dashboards (Product Mgmt)", "Self-Service (Enable End-Users)",
        "User Access Reqs & Ad-hoc Access", "Create Ad-Hoc Reports & Dashboards", "Data Wrangling",
        "Data Discovery", "Business Intelligence", "Data Science",
        "Data Policy & Standards", "Data Quality Mgmt", "Data Acquisition, Procurement, Integration",
        "Data Platform Engineering & Mgmt", "Data Infrastructure Mgmt", "Data Transformation & Integration",
        "Data Tool Lifecycle Mgmt & Support", "Data Integration", "Data Storage & Recovery",
        "Semantic Layer & Custom Subject Areas", "Data Privacy & Compliance", "Data Access Controls",
        "BI & Data Tool Access Mgmt", "Dashboard & Report Access Controls", "Trustworthy AI & Risk Mgmt",
        "AI/ML Solution Build (for Reporting)", "AI Model Governance",
    ],
}

PE = {
    "title": "Platform Engineering", "ncols": 2,
    "items": [
        "IoT Platform Mgmt", "IaaS Mgmt",
        "SaaS Mgmt", "API/Integration Mgmt",
        "Automation Platform Mgmt", "DevSecOps Platform Mgmt",
        "Cloud Platform Development", "Cloud Operations",
        "Cloud Financial Mgmt",
    ],
}

NE = {
    "title": "Network Engineering", "ncols": 2,
    "items": [
        "Network Operations", "WAN/LAN Mgmt",
        "Firewall Mgmt", "Telecom Mgmt",
        "Network Reliability Engineering",
    ],
}

PEM = {
    "title": "Partnerships & Ecosystem Management", "ncols": 4,
    "items": [
        "Business Relationships Mgmt", "Inter-IT Relationships", "Partner Strategy", "SLA Development",
        "SLA Metrics & Reporting", "Continuous Services Improvement", "Chargeback", "Service Quality Mgmt",
        "Tech Service Perf Mgmt", "Service Delivery Standards Mgmt", "Sourcing & Procurement Mgmt", "Vendor Mgmt",
        "Multi Service Provider Integration", "Supplier Relation Mgmt", "Supplier Performance Mgmt", "Contract Mgmt",
        "Supplier Relationship Mgmt",
    ],
}

ITES = {
    "title": "IT Enabling Services", "ncols": 4,
    "items": [
        "Contact Center", "Service Desk/Field Services", "Desktop Mgmt", "Collaboration Mgmt",
        "Unified Comms Mgmt", "Production Operations", "Technology Help Desk", "Incident Mgmt",
        "Problem Mgmt", "Tech Marketing & Comms",
    ],
}

TALENT = {
    "title": "Technology Talent", "ncols": 4,
    "items": [
        "Talent Strategy", "Workforce Strategy & Planning", "Talent Acquisition & Branding", "Succession Planning",
        "Performance & Rewards Mgmt", "Organizational Change Mgmt", "Learning & Development", "Career Mobility & Retention",
    ],
}

# ---------------------------------------------------------------------------
# LAYOUT GEOMETRY (inches) on a 13.333 x 7.5 widescreen slide
# ---------------------------------------------------------------------------
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.08

HEADER_H   = 0.16   # section header bar height
GROUP_H    = 0.20   # top-level group header height
BOX_H      = 0.225  # capability box height
VGAP       = 0.02   # vertical gap between box rows
HGAP       = 0.03   # horizontal gap between boxes
PANEL_GAP  = 0.08   # vertical gap between stacked panels
FONT_BOX   = 5.0    # box label font (pt)
FONT_HDR   = 6.5    # section header font
FONT_GRP   = 10.0   # group header font

prs = Presentation()
prs.slide_width  = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# running color index so the whole board cycles through the mix
_color_idx = [0]


def _set_text(tf_para_run, text, size, color, bold=False):
    run = tf_para_run
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = "Arial"


def add_box(x, y, w, h, text, fill, font_size, font_color,
            bold=False, line=BOX_LINE, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = line
    sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.02))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_text(p.add_run(), text, font_size, font_color, bold)
    return sp


def render_panel(panel, x, y, width):
    """Render a section panel (header bar + grid of boxes). Returns bottom y."""
    # header bar
    add_box(x, y, width, HEADER_H, panel["title"], HEADER_BG,
            FONT_HDR, WHITE, bold=True, line=HEADER_BG, line_w=0.5)
    y_cur = y + HEADER_H + VGAP
    ncols = panel["ncols"]
    bw = (width - (ncols + 1) * HGAP) / ncols
    items = panel["items"]
    for i, label in enumerate(items):
        r, c = divmod(i, ncols)
        bx = x + HGAP + c * (bw + HGAP)
        by = y_cur + r * (BOX_H + VGAP)
        add_box(bx, by, bw, BOX_H, label, color_for(_color_idx[0]),
                FONT_BOX, TEXT_DK)
        _color_idx[0] += 1
    nrows = (len(items) + ncols - 1) // ncols
    return y_cur + nrows * (BOX_H + VGAP)


def group_header(x, y, width, text):
    add_box(x, y, width, GROUP_H, text, HEADER_BG, FONT_GRP, WHITE,
            bold=True, line=HEADER_BG, line_w=0.5)
    return y + GROUP_H + VGAP


# --- Title -----------------------------------------------------------------
tb = slide.shapes.add_textbox(Inches(0), Inches(0.02), Inches(SLIDE_W), Inches(0.45))
p = tb.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
_set_text(p.add_run(), "AI's impact across the Tech Org: Capability View",
          18, TEXT_DK, bold=True)

CONTENT_TOP = 0.55

# --- LEFT: Product group ----------------------------------------------------
prod_x, prod_w = MARGIN, 2.05
y = group_header(prod_x, CONTENT_TOP, prod_w, "Product")
render_panel(PRODUCT_PANEL, prod_x, y, prod_w)

# --- RIGHT: Tech Platforms, Services, Operations group ----------------------
right_x = prod_x + prod_w + 0.10
right_w = SLIDE_W - right_x - MARGIN
gy = group_header(right_x, CONTENT_TOP, right_w, "Tech Platforms, Services, Operations")

# three masonry columns inside the right group
colA_x, colA_w = right_x,                 3.55
colB_x, colB_w = colA_x + colA_w + 0.10,  2.95
colC_x, colC_w = colB_x + colB_w + 0.10,  right_x + right_w - (colB_x + colB_w + 0.10)

# Column A
yA = render_panel(TBM, colA_x, gy, colA_w) + PANEL_GAP
yA = render_panel(EA,  colA_x, yA, colA_w) + PANEL_GAP
yA = render_panel(DA,  colA_x, yA, colA_w)

# Column B
yB = render_panel(TSI, colB_x, gy, colB_w) + PANEL_GAP
yB = render_panel(IE,  colB_x, yB, colB_w) + PANEL_GAP
yB = render_panel(PE,  colB_x, yB, colB_w) + PANEL_GAP
yB = render_panel(NE,  colB_x, yB, colB_w)

# Column C
yC = render_panel(ISM, colC_x, gy, colC_w) + PANEL_GAP
yC = render_panel(PEM, colC_x, yC, colC_w) + PANEL_GAP
yC = render_panel(ITES, colC_x, yC, colC_w)

# Technology Talent: full-width band across the right group, below columns
talent_y = max(yA, yB, yC) + PANEL_GAP
render_panel(TALENT, right_x, talent_y, right_w)

# --- Legend -----------------------------------------------------------------
talent_bottom = talent_y + GROUP_H  # talent panel header + small grid bottom approx
leg_y = max(talent_y + 0.95, 6.85)
leg_y = min(leg_y, SLIDE_H - 0.35)
legend = [
    ("= AI Powered\n80% AI / 20% Human",   POWERED),
    ("= AI Augmented\n50% AI / 50% Human",  AUGMENTED),
    ("= AI Assisted\n20% AI / 80% Human",   ASSISTED),
]
# legend title
lt = slide.shapes.add_textbox(Inches(right_x), Inches(leg_y - 0.02),
                              Inches(3.3), Inches(0.3))
lt.text_frame.word_wrap = True
lp = lt.text_frame.paragraphs[0]
lp.alignment = PP_ALIGN.LEFT
_set_text(lp.add_run(), "LEGEND: ", 11, TEXT_DK, bold=True)
r_it = lp.add_run(); _set_text(r_it, "Degree of Expected AI Impact", 11, TEXT_DK, bold=True)
r_it.font.italic = True

lx = right_x + 3.3
step = (SLIDE_W - MARGIN - lx) / 3.0   # evenly distribute across remaining width
for text, color in legend:
    swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(lx), Inches(leg_y), Inches(0.26), Inches(0.26))
    swatch.fill.solid(); swatch.fill.fore_color.rgb = color
    swatch.line.color.rgb = RGBColor(0x80, 0x80, 0x80); swatch.line.width = Pt(0.5)
    swatch.shadow.inherit = False
    tb2 = slide.shapes.add_textbox(Inches(lx + 0.30), Inches(leg_y - 0.06),
                                   Inches(step - 0.34), Inches(0.5))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    pp = tf2.paragraphs[0]
    _set_text(pp.add_run(), text.split("\n")[0], 8.5, TEXT_DK, bold=True)
    p2 = tf2.add_paragraph()
    _set_text(p2.add_run(), text.split("\n")[1], 7.5, TEXT_DK, bold=False)
    p2.runs[0].font.italic = True
    lx += step

prs.save("AI_Capability_View.pptx")
print("Saved AI_Capability_View.pptx")
